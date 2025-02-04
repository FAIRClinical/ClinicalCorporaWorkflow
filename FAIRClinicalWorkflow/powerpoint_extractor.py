from pptx import Presentation
presentation_extensions = [".pptx", ".ppt", ".pptm", ".odp"]

def get_powerpoint_text(filename):
    text = []
    pres = Presentation(filename)
    for slide in pres.slides:
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                text.append(shape.text)
    return text
